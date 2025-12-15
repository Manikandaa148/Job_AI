from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    # Create a presentation object
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Job AI: Revolutionizing Job Search with AI"
    subtitle.text = "A Comprehensive Auto-Apply & Job Management System\nPresented by: Team Job AI\nDecember 2025"

    # Helper to add bullet slides
    def add_bullet_slide(heading, points):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = heading
        tf = slide.placeholders[1].text_frame
        tf.text = points[0] # First point
        for point in points[1:]:
            p = tf.add_paragraph()
            p.text = point

    # --- Slide 2: The Problem ---
    add_bullet_slide("The Problem: Job Hunting is Broken", [
        "Time-Consuming: Applicants spend hours filling out repetitive forms.",
        "Complexity: Multiple platforms, different requirements, and varying formats.",
        "Lack of Feedback: 'Apply and forget' with no tracking.",
        "Profile Management: Keeping data updated across sites is a hassle."
    ])

    # --- Slide 3: The Solution ---
    add_bullet_slide("The Solution: Job AI", [
        "What is Job AI? An intelligent, all-in-one platform for automation.",
        "Core Promise: 'One-Click Application' with intelligent validation.",
        "Differentiation: Active agents working for you, not just a job board.",
        "Result: Saves time, increases volume, improves quality."
    ])

    # --- Slide 4: Key Features ---
    add_bullet_slide("Key Features Overview", [
        "⚡ Intelligent Auto-Apply: One-click submission with AI validation.",
        "🤖 Interactive Chatbot: Real-time assistant for profile completion.",
        "🔍 Advanced Job Search: Aggregated jobs from multiple sources.",
        "📄 Smart Profile Management: Centralized profile adapting to forms.",
        "🎨 Premium UI/UX: Dark-mode enabled, modern interface."
    ])

    # --- Slide 5: Feature Spotlight: Auto-Apply ---
    add_bullet_slide("Spotlight: Intelligent Auto-Apply", [
        "Profile Validation Engine: Checks 8+ critical fields before sending.",
        "Missing Field Detection: Identifies exactly what is missing.",
        "Visual Status Tracking: Idle, Loading, Success, Error, Missing Info.",
        "Error Recovery: Gracefully handles submission errors."
    ])

    # --- Slide 6: Feature Spotlight: Chatbot ---
    add_bullet_slide("Spotlight: Interactive Chatbot", [
        "Context-Aware: Knows why an application paused.",
        "Proactive: Opens automatically when data is missing.",
        "Real-Time Updates: Chat input immediately updates the profile.",
        "User Engagement: Friendly conversation with history and typing indicators."
    ])

    # --- Slide 7: Technical Architecture ---
    add_bullet_slide("Technical Architecture", [
        "Frontend: Next.js 14, TypeScript, Tailwind CSS.",
        "Backend: FastAPI (Python), SQLAlchemy ORM.",
        "Database: SQLite (Dev) / PostgreSQL (Prod).",
        "Security: JWT Auth, Pydantic Validation, BCRYPT.",
        "AI Logic: Custom Python agents for validation."
    ])

    # --- Slide 8: UX & Design ---
    add_bullet_slide("User Experience (UX) & Design", [
        "Visual Aesthetics: Glassmorphism, Smooth Gradients, Micro-animations.",
        "Responsiveness: Fully functional on Mobile, Tablet, and Desktop.",
        "Accessibility: High contrast Dark Mode, clear feedback.",
        "Interactive Elements: Hover effects, floating buttons (Chatbot)."
    ])

    # --- Slide 9: User Journey ---
    add_bullet_slide("The User Journey", [
        "1. Search: User finds a relevant job card.",
        "2. Action: Clicks 'Auto Apply' (Purple Lightning button).",
        "3. Process: System validates and applies (or Chatbot helps if info missing).",
        "4. Result: User receives visual confirmation instantly."
    ])

    # --- Slide 10: Roadmap ---
    add_bullet_slide("Future Roadmap & Conclusion", [
        "Phase 2: Bulk Auto-Apply (50+ jobs).",
        "Phase 3: AI Cover Letters & Interview Tips.",
        "Phase 4: Direct Integrations (LinkedIn/Glassdoor API).",
        "Conclusion: Job AI is the future of automated recruitment."
    ])

    # Save
    output_path = os.path.join(os.getcwd(), 'Job_AI_Presentation.pptx')
    prs.save(output_path)
    print(f"Presentation saved successfully at: {output_path}")

if __name__ == "__main__":
    create_presentation()
